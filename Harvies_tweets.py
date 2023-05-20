# -*- coding: utf-8 -*-
"""
Created on Sat May 22 10:22:26 2021

@author: DELL
"""

from pptx import Presentation
import tweepy
import pandas as pd
import csv
import numpy as np


##funcion para saber si es retweet o tweet normal
def isRetweetd(tweet_txt=""):
    rt="RT"
    contador=0
    contador_tweet=0
    bandera=False
    for letra in tweet_txt:
        if(letra==rt[0] or letra==rt[1]):
            contador_tweet+=1
        if(contador_tweet==2):
            bandera=True
        if(contador>3):
            return bandera
        contador+=1
    print(bandera)
    return bandera



##Claves
consumer_key = '4PtaaR1vl8Zfb8ZXRXrzAcbEm' 
consumer_secret = 'DhXvnqopuqd98yBZ0uTSi5WhKmX0GpeOK5T1mPUkUZrZEZfrKt'
#Autenticacion
auth = tweepy.OAuthHandler(consumer_key, consumer_secret)
#API
api = tweepy.API(auth)
#VAriables TWEETS(URL,)
tweets = api.user_timeline(screen_name='FJG_TD',tweetmode="extended")

user = api.get_user(screen_name='FJG_TD')
tweets_url=[]
tweets_date=[]
tweets_likes=[]
tweets_retweets=[]
tweets_respuestas=[]##No lo tengo
tweets_Seguidores=[]
tweets_ti=[]
tweet_texto=[]

tweet_reply=[]
nada='nada'
lista=[tweets_retweets,tweet_texto]

#Ciclo de datos a meter a excel
for tweet in tweets:
  #Fucion para saber si es retweet
  if(not isRetweetd(tweet_txt=tweet.text)):
         k=tweet.entities['urls']
         print(k)
         #BUg url
         if not k:
             tweets_url.append(nada)
         else:
             for twe in k:
                 tweets_url.append(twe['expanded_url'])
                 
        #es, , variables
         tweets_retweets.append(tweet.retweet_count)
         tweets_date.append(tweet.created_at)
         tweet_texto.append(tweet.text)
         tweets_likes.append(tweet.favorite_count)
         tweets_Seguidores.append(user.followers_count)
         
#Dtaframe para importar con la libreria

df=pd.DataFrame({'a':tweets_url,'b':tweets_date,'c':tweets_likes,'d':tweets_retweets
                 ,'e':tweets_Seguidores})


   
#Guardmos en 
export_csv=df.to_csv(r'C:\Users\DELL\Documents\BaseDatosSuccces\reciente.csv')

print(tweets[0])


#Presentacion
prs = Presentation(r'C:\Users  \DELL\Documents\BaseDatosSuccces\otros2.pptx')
title_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Cambio!"
subtitle.text = "python-pptx was here!"

prs.save(r'C:\Users\DELL\Documents\BaseDatosSuccces\otros2.pptx')